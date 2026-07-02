import cv2
from flask import jsonify
from flask import Flask, render_template, request, redirect, session
from authlib.integrations.flask_client import OAuth
import sqlite3
import os
from werkzeug.utils import secure_filename
from pptx import Presentation
import pdfplumber
import google.generativeai as genai
import json
from werkzeug.utils import secure_filename
import os
from video_analysis import analyze_video
from dotenv import load_dotenv
load_dotenv()

app = Flask(__name__)
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
genai.configure(api_key=GEMINI_API_KEY)

model = genai.GenerativeModel("gemini-2.5-flash")
app.secret_key = "presentiq_secret_key"
oauth = OAuth(app)


UPLOAD_FOLDER = "uploads"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)


google = oauth.register(
    name="google",
    client_id=os.getenv("GOOGLE_CLIENT_ID"),
    client_secret=os.getenv("GOOGLE_CLIENT_SECRET"),
    server_metadata_url="https://accounts.google.com/.well-known/openid-configuration",
    client_kwargs={
        "scope": "openid email profile"
    }
)

# ---------------- DATABASE ---------------- #

def init_db():
    conn = sqlite3.connect("database.db")
    c = conn.cursor()

    c.execute("""
    CREATE TABLE IF NOT EXISTS users(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT,
        email TEXT UNIQUE,
        password TEXT
    )
    """)

    conn.commit()
    conn.close()

init_db()

# ---------------- HOME ---------------- #

@app.route("/")
def home():
    return render_template("index.html")


# ---------------- SIGNUP ---------------- #

@app.route("/signup", methods=["GET", "POST"])
def signup():

    if request.method == "POST":

        name = request.form["name"]
        email = request.form["email"]
        password = request.form["password"]

        conn = sqlite3.connect("database.db")
        c = conn.cursor()

        try:
            c.execute(
                "INSERT INTO users(name,email,password) VALUES(?,?,?)",
                (name, email, password)
            )

            conn.commit()
            conn.close()

            return redirect("/login")

        except Exception as e:
            conn.close()
            return str(e)

    return render_template("signup.html")

@app.route("/google")
def google_login():
    redirect_uri = "http://127.0.0.1:5000/auth"
    return google.authorize_redirect(redirect_uri)


@app.route("/auth")
def auth():
    token = google.authorize_access_token()
    user = token["userinfo"]

    session["username"] = user["name"]

    return redirect("/dashboard")


# ---------------- LOGIN ---------------- #

@app.route("/login", methods=["GET", "POST"])
def login():

    if request.method == "POST":

        email = request.form["email"]
        password = request.form["password"]

        conn = sqlite3.connect("database.db")
        c = conn.cursor()

        c.execute(
            "SELECT name FROM users WHERE email=? AND password=?",
            (email, password)
        )

        user = c.fetchone()
        conn.close()

        if user:
            session["username"] = user[0]
            session["email"] = email
            return redirect("/dashboard")

        else:
            return render_template(
                "login.html",
                error="Invalid Email or Password"
            )

    return render_template("login.html")


# ---------------- DASHBOARD ---------------- #

@app.route("/dashboard")
def dashboard():

    if "username" not in session:
        return redirect("/login")

    return render_template(
        "dashboard.html",
        username=session["username"],
        email=session["email"]
    )

@app.route("/upload", methods=["GET", "POST"])
def upload():

    if request.method == "POST":

        if "ppt" not in request.files:
            return "No file selected"

        file = request.files["ppt"]

        if file.filename == "":
            return "Please choose a PPT"

        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)

        file.save(filepath)
        if filename.lower().endswith(".pdf"):
            ppt_text = extract_pdf_text(filepath)
        else:
            ppt_text = extract_ppt_text(filepath)
        ai_result = analyze_with_gemini(ppt_text)
        
        session["ppt_text"] = ppt_text
        session["ppt_path"] = filepath

        return render_template(
            "analysis.html",
            analysis=ai_result
        )

    return render_template("upload.html")


def extract_ppt_text(filepath):
    prs = Presentation(filepath)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text
def extract_pdf_text(filepath):

    text = ""

    with pdfplumber.open(filepath) as pdf:

        for page in pdf.pages:

            page_text = page.extract_text()

            if page_text:
                text += page_text + "\n"

    return text
def analyze_with_gemini(text):
    prompt = f"""
    You are an expert AI Presentation Coach.

    Analyze the following presentation professionally.

    Return ONLY plain text.

    Use this exact format.

    🤖 AI Presentation Analysis

    🎯 Overall Score
    Score: XX/100

    Write 2-3 sentences explaining why the presentation received this score.

    📚 Content Quality
    Score: X/10

    Explain the quality of the content in one short paragraph.

    🔄 Presentation Flow
    Score: X/10

    Explain whether the slides follow a logical order.

    🎨 Visual Design
    Score: X/10

    Evaluate the slide design based on the presentation text. Mention if diagrams, images, charts or visual elements seem to be missing.

    ✅ Strengths

    Give exactly 5 bullet points.

    ⚠️ Weaknesses

    Give exactly 5 bullet points.

    💡 Suggestions for Improvement

    Give exactly 5 numbered suggestions that are practical and specific.

    📝 Final Summary

    Write a professional conclusion in 3-4 sentences.

    Rules:
    - Do NOT use Markdown.
    - Do NOT use #.
    - Do NOT use *.
    - Do NOT use **.
    - Do NOT use tables.
    - Do NOT use code blocks.
    - Do NOT use horizontal lines.
    - Use only the icons already provided in the headings.
    - Always return scores as:
    Score: XX/100
    Score: X/10
    - Keep the report professional, neat and easy to read.
    - Do not repeat the same points.

    Presentation Text:

    {text}
    """

    response = model.generate_content(prompt)
    return response.text
@app.route("/live_practice")
def live_practice():

    if "username" not in session:
        return redirect("/login")

    return render_template("live_practice.html")
from flask import request, jsonify

@app.route("/analyze_speech", methods=["POST"])
def analyze_speech():

    data = request.get_json()

    speech = data.get("speech", "")

    prompt = f"""
    You are an AI Presentation Coach.

    Analyze this presentation.

    Presentation:
    {speech}

    Return ONLY JSON.

    {{
    "score":"",
    "confidence":"",
    "nervousness":"",
    "communication":"",
    "strengths":[
    "",
    "",
    ""
    ],
    "improvements":[
    "",
    "",
    ""
    ],
    "feedback":""
    }}

    Do not return Markdown.
    Do not explain anything.
    Return valid JSON only.
    """
    response = model.generate_content(prompt)

    result = json.loads(response.text)

    return jsonify(result)
  

@app.route("/upload_video", methods=["GET", "POST"])
def upload_video():

    if request.method == "POST":

        video = request.files["video"]

        filename = secure_filename(video.filename)

        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)

        video.save(filepath)

        result = analyze_video(filepath)

        return render_template("reports.html", report=result)

    return render_template("upload_video.html")
@app.route("/reports")
def reports():
    return render_template("reports.html")

# ---------------- LOGOUT ---------------- #

@app.route("/logout")
def logout():
    session.clear()
    return redirect("/")


# ---------------- RUN ---------------- #

if __name__ == "__main__":
    app.run(debug=True)